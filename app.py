import collections.abc
import config
assert collections
import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI
from io import BytesIO
import requests
from PIL import Image

# openai client
openai_client = OpenAI(
    api_key=config.OPENAI_API_KEY
)


def slide_generator(text, prs):
    # 创建生成 DALL-E 图像的提示
    image_prompt = f"Summarize the following text to a DALL-E image generation prompt: \n{text}"
    model_engine = "gpt-4o"

    dlp = openai_client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{image_prompt}"}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )

    dalle_prompt = dlp.choices[0].message.content

    model_img = "dall-e-3"
    response = openai_client.images.generate(
        model=model_img,
        prompt=dalle_prompt + " Style: digital art",
        n=1,
        size="1024x1024"
    )
    image_url = response.data[0].url

    # 下载生成的图像
    image_response = requests.get(image_url)
    image_stream = BytesIO(image_response.content)
    image = Image.open(image_stream)

    # 将图像保存为临时文件
    image_path = "temp_image.png"
    image.save(image_path)

    # 创建生成 PowerPoint 要点文本的提示
    bullet_prompt = f"Create a bullet point text for a Powerpoint slide from the following text: \n{text}"
    ppt_bullet = openai_client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{bullet_prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = ppt_bullet.choices[0].message.content

    # 创建生成 PowerPoint 标题的提示
    title_prompt = f"Create a title for a Powerpoint slide from the following text: \n{text}"
    ppt_title = openai_client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{title_prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = ppt_title.choices[0].message.content

    # 将图像和文本添加到幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = ppt_header
    body.text = ppt_text
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(6))

def get_slides():
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)
    prs.save("my_presentation.pptx")

app = tk.Tk()
app.title("Create PPT Slides")
app.geometry("800x600")

# 创建文本字段
text_field = tk.Text(app)
text_field.pack(fill="both", expand=True)
text_field.configure(wrap="word", font=("Arial", 12))
text_field.focus_set()

# 创建生成幻灯片的按钮
create_button = tk.Button(app, text="Create Slides", command=get_slides)
create_button.pack()

app.mainloop()
